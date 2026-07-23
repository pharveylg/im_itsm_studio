#!/usr/bin/env python3
"""Generate the OAS Major Incident showcase deck as a native .pptx (16:9)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# ---- palette ----
NAVY   = RGBColor(0x0F, 0x17, 0x2A)
PANEL  = RGBColor(0x1E, 0x29, 0x3B)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
GREEN  = RGBColor(0x10, 0xB9, 0x81)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xE2, 0xE8, 0xF0)
BLUET  = RGBColor(0xBF, 0xDB, 0xFE)
PURPT  = RGBColor(0xC4, 0xB5, 0xFD)
GREY   = RGBColor(0x64, 0x74, 0x8B)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SW, SH = prs.slide_width, prs.slide_height

# ---------- low-level helpers ----------
def set_bg(slide, color=NAVY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def _txt(tf, text, size, color, bold=False, align=PP_ALIGN.CENTER, font='Calibri'):
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font

def box(slide, x, y, w, h, text, fill=BLUE, tcolor=WHITE, size=11, bold=False,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=None, align=PP_ALIGN.CENTER):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    tf = sp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(5); tf.margin_right = Pt(5)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    _txt(tf, text, size, tcolor, bold, align)
    return sp

def arrow(slide, x1, y1, x2, y2, color=BLUE, dashed=False, width=1.5):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color; conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    tail = ln.find(qn('a:tailEnd'))
    if tail is None:
        tail = ln.makeelement(qn('a:tailEnd'), {}); ln.append(tail)
    tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')
    if dashed:
        ln.set('prstDash', 'dash')
    return conn

def connect(slide, a, b, color=BLUE, dashed=False):
    ax, ay, aw, ah = a.left, a.top, a.width, a.height
    bx, by, bw, bh = b.left, b.top, b.width, b.height
    acx, acy = ax + aw//2, ay + ah//2
    bcx, bcy = bx + bw//2, by + bh//2
    if abs(bcx - acx) >= abs(bcy - acy):
        if bcx >= acx:
            x1, y1, x2, y2 = ax+aw, acy, bx, bcy
        else:
            x1, y1, x2, y2 = ax, acy, bx+bw, bcy
    else:
        if bcy >= acy:
            x1, y1, x2, y2 = acx, ay+ah, bcx, by
        else:
            x1, y1, x2, y2 = acx, ay, bcx, by+bh
    arrow(slide, x1, y1, x2, y2, color, dashed)

def new_slide():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    return s

def chrome(slide):
    b = slide.shapes.add_textbox(Inches(0.4), Inches(0.06), Inches(6), Inches(0.3))
    _txt(b.text_frame, "OAS · Operational Analysis Standard", 10, BLUET, True, PP_ALIGN.LEFT)
    f = slide.shapes.add_textbox(Inches(8.5), Inches(7.05), Inches(4.4), Inches(0.3))
    _txt(f.text_frame, "Major Incident Showcase · Release 1.2.2", 9, GREY, False, PP_ALIGN.RIGHT)

def title(slide, text, sub=None):
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.8))
    _txt(t.text_frame, text, 26, BLUET, True, PP_ALIGN.LEFT)
    if sub:
        s = slide.shapes.add_textbox(Inches(0.5), Inches(1.12), Inches(12.3), Inches(0.4))
        _txt(s.text_frame, sub, 13, LIGHT, False, PP_ALIGN.LEFT)

def bullets(slide, x, y, w, h, items, size=14, color=LIGHT, gap=6):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = "•  " + it
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = 'Calibri'
    return tb

def card(slide, x, y, w, h, text, accent=BLUE, tcolor=LIGHT, size=12, bold=False):
    c = box(slide, x, y, w, h, text, fill=PANEL, tcolor=tcolor, size=size, bold=bold,
            line=accent, align=PP_ALIGN.LEFT)
    return c

def table(slide, x, y, w, h, rows, size=11, col0=4.0):
    nr, nc = len(rows), len(rows[0])
    gf = slide.shapes.add_table(nr, nc, x, y, w, h)
    tbl = gf.table
    # proportional widths
    avail = w
    first = Inches(col0)
    tbl.columns[0].width = first
    rest = (avail - first) // (nc - 1)
    for c in range(1, nc):
        tbl.columns[c].width = rest
    for r in range(nr):
        for c in range(nc):
            cell = tbl.cell(r, c)
            cell.text = str(rows[r][c])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Pt(5); cell.margin_right = Pt(5)
            cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(size); p.font.name = 'Calibri'
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1E,0x3A,0x8A)
                p.font.color.rgb = BLUET; p.font.bold = True
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
                p.font.color.rgb = LIGHT
    return tbl

def tag(slide, x, y, text, color=BLUE):
    return box(slide, x, y, Inches(len(text)*0.105+0.5), Inches(0.34), text,
               fill=RGBColor(0x16,0x29,0x4A), tcolor=BLUET, size=10, bold=True,
               line=color, align=PP_ALIGN.CENTER)

# ================= SLIDES =================

# 1 — Title
s = new_slide(); chrome(s)
t = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.2))
_txt(t.text_frame, "Operational Analysis Standard", 40, WHITE, True, PP_ALIGN.CENTER)
st = s.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.7))
_txt(st.text_frame, "Major Incident Analysis in Context", 24, BLUET, True, PP_ALIGN.CENTER)
sl = s.shapes.add_textbox(Inches(1.5), Inches(4.4), Inches(10.3), Inches(0.6))
_txt(sl.text_frame, "How OAS-201 Major Incident Communications connects to every other standard", 14, LIGHT, False, PP_ALIGN.CENTER)
tags = ["Release 1.2.2", "ITIL 4-aligned", "ServiceNow reference", "Evidence-based"]
tx = Inches(3.0)
for tg in tags:
    tag(s, tx, Inches(5.4), tg, BLUE); tx += Inches(len(tg)*0.105+0.7)

# 2 — What is OAS
s = new_slide(); chrome(s); title(s, "What is OAS? (through the Major Incident lens)")
card(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(1.0),
     "A governed library of evidence-based analysis standards for operational records.", BLUE)
card(s, Inches(0.5), Inches(2.9), Inches(6.0), Inches(1.0),
     "It complements ITIL & ServiceNow — it does not replace them.", PURPLE)
card(s, Inches(0.5), Inches(4.1), Inches(6.0), Inches(1.0),
     "Every standard inherits OAS-000 and the evidence-first principle.", GREEN)
card(s, Inches(6.9), Inches(1.7), Inches(5.9), Inches(1.0),
     "Today's focus: OAS-201 — Major Incident Communications.", BLUE)
card(s, Inches(6.9), Inches(2.9), Inches(5.9), Inches(1.0),
     "OAS-201 evaluates HOW an incident was communicated — not the fix.", PURPLE)
card(s, Inches(6.9), Inches(4.1), Inches(5.9), Inches(1.0),
     "…and how it connects to OAS-000, 101, 301, 401, 501 and the KB.", GREEN)

# 3 — Library at a glance (architecture)
s = new_slide(); chrome(s); title(s, "The OAS Library at a glance",
                                          "OAS-000 governs all; the Knowledge Base is derived")
# Governance col
box(s, Inches(0.5), Inches(1.45), Inches(2.7), Inches(0.4), "Governance", fill=PANEL, tcolor=BLUET, size=12, bold=True)
g0 = box(s, Inches(0.5), Inches(2.0), Inches(2.7), Inches(0.7), "OAS-000\nGovernance", fill=BLUE)
g1 = box(s, Inches(0.5), Inches(3.0), Inches(2.7), Inches(0.7), "OAS-501\nKnowledge", fill=BLUE)
# Methodologies col
box(s, Inches(5.0), Inches(1.45), Inches(2.7), Inches(0.4), "Methodologies", fill=PANEL, tcolor=BLUET, size=12, bold=True)
m1 = box(s, Inches(5.0), Inches(1.95), Inches(2.7), Inches(0.6), "OAS-101 Incident", fill=BLUE)
m2 = box(s, Inches(5.0), Inches(2.65), Inches(2.7), Inches(0.6), "OAS-201 MI Comms", fill=PURPLE)
m3 = box(s, Inches(5.0), Inches(3.35), Inches(2.7), Inches(0.6), "OAS-301 Problem", fill=BLUE)
m4 = box(s, Inches(5.0), Inches(4.05), Inches(2.7), Inches(0.6), "OAS-401 Change", fill=BLUE)
# KB col
box(s, Inches(9.6), Inches(1.45), Inches(3.0), Inches(0.4), "Knowledge Base", fill=PANEL, tcolor=PURPT, size=12, bold=True)
kb1 = box(s, Inches(9.6), Inches(2.0), Inches(3.0), Inches(0.9), "KB-001 … KB-006", fill=PURPLE, tcolor=WHITE)
kb2 = box(s, Inches(9.6), Inches(3.1), Inches(3.0), Inches(0.9), "Templates … Prompts", fill=PURPLE, tcolor=WHITE)
# arrows
for m in (m1, m2, m3, m4):
    connect(s, g0, m, BLUE, dashed=True)
connect(s, g0, g1, BLUE)
for m in (m1, m2, m3, m4):
    connect(s, m, kb1, PURPLE)
connect(s, g1, kb1, PURPLE)
# legend
box(s, Inches(0.5), Inches(5.4), Inches(6.0), Inches(0.4),
    "Dashed blue = OAS-000 governs   ·   Purple = Knowledge Base derived", fill=PANEL,
    tcolor=GREY, size=10, align=PP_ALIGN.LEFT)

# 4 — Where MI lives
s = new_slide(); chrome(s); title(s, "Where Major Incident lives")
inc = box(s, Inches(0.6), Inches(3.1), Inches(2.6), Inches(0.9), "Incident", fill=BLUE)
mi  = box(s, Inches(4.0), Inches(3.1), Inches(2.6), Inches(0.9), "Major Incident", fill=PURPLE)
chg = box(s, Inches(8.2), Inches(1.6), Inches(2.6), Inches(0.8), "Remediation Change", fill=BLUE)
prb = box(s, Inches(8.2), Inches(3.1), Inches(2.6), Inches(0.8), "Problem", fill=BLUE)
ka  = box(s, Inches(8.2), Inches(4.6), Inches(2.6), Inches(0.8), "Knowledge", fill=PURPLE)
connect(s, inc, mi, BLUE)
connect(s, mi, chg, BLUE)
connect(s, mi, prb, BLUE)
connect(s, mi, ka, PURPLE)
box(s, Inches(0.6), Inches(5.4), Inches(11.0), Inches(0.9),
    "A Major Incident is declared from an Incident, frequently drives a Problem and remediation "
    "Changes, and generates operational knowledge.", fill=PANEL, tcolor=LIGHT, size=12, align=PP_ALIGN.LEFT)

# 5 — OAS-201 focus
s = new_slide(); chrome(s); title(s, "OAS-201 — Major Incident Communications (the focus)")
card(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(1.1),
     "Purpose: evaluate effectiveness, timeliness, consistency, clarity & governance of communications.", BLUE)
card(s, Inches(0.5), Inches(3.0), Inches(6.0), Inches(1.1),
     "Scope: operational updates, executive, customer, vendor, bridge, MIM handovers.", PURPLE)
card(s, Inches(6.9), Inches(1.7), Inches(5.9), Inches(1.1),
     "Does: assess HOW the incident was communicated.", BLUE)
card(s, Inches(6.9), Inches(3.0), Inches(5.9), Inches(1.1),
     "Does NOT: re-do technical root-cause (OAS-301) or incident response (OAS-101).", GREEN)

# 6 — OAS-201 phases
s = new_slide(); chrome(s); title(s, "OAS-201 — 8-phase methodology")
ph = ["1 Context","2 Comms Timeline","3 Quality","4 Timeliness",
      "5 Stakeholders","6 MIM Handovers","7 Governance","8 Effectiveness"]
xs = [0.5, 3.55, 6.6, 9.65]
ys = [Inches(2.2), Inches(4.3)]
boxes = []
for i, p in enumerate(ph):
    col = i % 4; row = i // 4
    b = box(s, Inches(xs[col]), ys[row], Inches(2.5), Inches(0.9), p, fill=BLUE)
    boxes.append(b)
for i in range(3):
    connect(s, boxes[i], boxes[i+1], BLUE)
connect(s, boxes[3], boxes[4], BLUE)
for i in range(4, 7):
    connect(s, boxes[i], boxes[i+1], BLUE)

# 7 — Relationship map
s = new_slide(); chrome(s); title(s, "Relationship map — Major Incident at the centre")
center = box(s, Inches(4.9), Inches(3.1), Inches(3.0), Inches(1.0), "OAS-201\nMajor Incident Comms", fill=PURPLE, tcolor=WHITE, size=12, bold=True)
o0 = box(s, Inches(5.2), Inches(1.2), Inches(3.0), Inches(0.8), "OAS-000 Governance", fill=BLUE)
o1 = box(s, Inches(0.5), Inches(3.0), Inches(2.6), Inches(0.8), "OAS-101 Incident", fill=BLUE)
o3 = box(s, Inches(0.5), Inches(5.0), Inches(2.6), Inches(0.8), "OAS-301 Problem", fill=BLUE)
o4 = box(s, Inches(4.3), Inches(5.3), Inches(2.6), Inches(0.8), "OAS-401 Change", fill=BLUE)
o5 = box(s, Inches(10.1), Inches(3.0), Inches(2.6), Inches(0.8), "OAS-501 Knowledge", fill=BLUE)
kb = box(s, Inches(10.1), Inches(1.2), Inches(2.6), Inches(0.8), "Knowledge Base", fill=PURPLE)
for tgt in (o0, o1, o3, o4, o5, kb):
    connect(s, center, tgt, BLUE)
connect(s, o0, center, BLUE, dashed=True)

# 8 — OAS-000
s = new_slide(); chrome(s); title(s, "↔ OAS-000 — Governance (the foundation)")
card(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(1.0),
     "Inherits: Evidence First · Objectivity · Traceability · Transparency · Fact/Opinion separation.", BLUE)
card(s, Inches(0.5), Inches(2.9), Inches(6.0), Inches(1.0),
     "Uses: Evidence Hierarchy (7 tiers), States, Classification, Confidence Model (H/M/L/U).", PURPLE)
card(s, Inches(6.9), Inches(1.7), Inches(5.9), Inches(1.0),
     "Lifecycle: Collect → Inventory → Classify → Analyse → Correlate → Timeline → Gaps → Confidence → Assess → Summarise.", GREEN)
card(s, Inches(6.9), Inches(2.9), Inches(5.9), Inches(1.0),
     "Rule: no methodology may contradict OAS-000.", BLUE)

# 9 — OAS-101 table
s = new_slide(); chrome(s); title(s, "↔ OAS-101 — Incident Analysis")
table(s, Inches(0.6), Inches(1.8), Inches(12.1), Inches(3.0),
      [["Dimension","OAS-101 (Incident)","OAS-201 (Major Incident Comms)"],
       ["Question","What happened & how was it handled?","How was it communicated?"],
       ["Output feeds","Operational facts, timeline, response","Stakeholder narrative & cadence"],
       ["Link","Primary record","Assesses comms about the same event"]])
card(s, Inches(0.6), Inches(5.1), Inches(12.1), Inches(1.0),
     "OAS-101 establishes the FACTS; OAS-201 assesses how those facts were COMMUNICATED. "
     "Their Related-Record assessment cross-links them.", BLUE)

# 10 — OAS-301
s = new_slide(); chrome(s); title(s, "↔ OAS-301 — Problem Analysis")
card(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(1.0),
     "A Major Incident often spawns a Problem (PRB).", BLUE)
card(s, Inches(0.5), Inches(2.9), Inches(6.0), Inches(1.0),
     "OAS-201's SITREPs & narrative feed OAS-301's investigation context.", PURPLE)
card(s, Inches(6.9), Inches(1.7), Inches(5.9), Inches(1.0),
     "OAS-301 validates root cause; OAS-201 confirms stakeholders were informed during.", GREEN)
card(s, Inches(6.9), Inches(2.9), Inches(5.9), Inches(1.0),
     "A Post Incident Review (PIR) ties both together.", BLUE)

# 11 — OAS-401
s = new_slide(); chrome(s); title(s, "↔ OAS-401 — Change Analysis")
card(s, Inches(0.5), Inches(1.7), Inches(12.1), Inches(0.9),
     "Remediation Changes (401B) frequently follow a Major Incident.", BLUE)
card(s, Inches(0.5), Inches(2.9), Inches(6.0), Inches(1.1),
     "OAS-401 assesses planning, risk, rollback symmetry & validation.", PURPLE)
card(s, Inches(6.9), Inches(2.9), Inches(5.9), Inches(1.1),
     "OAS-201 assesses comms during/after the change.", GREEN)
card(s, Inches(0.5), Inches(4.3), Inches(12.1), Inches(1.0),
     "Consistency check: did communications reflect the change's ACTUAL outcome (not just its closure)?", BLUE)

# 12 — OAS-501 diagram
s = new_slide(); chrome(s); title(s, "↔ OAS-501 — Operational Knowledge")
mi = box(s, Inches(0.6), Inches(3.0), Inches(2.6), Inches(0.9), "Major Incident\nLessons", fill=PURPLE)
ka = box(s, Inches(4.2), Inches(1.7), Inches(2.6), Inches(0.7), "Knowledge Article", fill=BLUE)
rb = box(s, Inches(4.2), Inches(3.0), Inches(2.6), Inches(0.7), "Runbook", fill=BLUE)
ke = box(s, Inches(4.2), Inches(4.3), Inches(2.6), Inches(0.7), "Known Error Article", fill=BLUE)
o5 = box(s, Inches(8.3), Inches(3.0), Inches(2.6), Inches(0.9), "OAS-501\ngoverns quality", fill=BLUE)
connect(s, mi, ka, PURPLE); connect(s, mi, rb, PURPLE); connect(s, mi, ke, PURPLE)
connect(s, ka, o5, BLUE); connect(s, rb, o5, BLUE); connect(s, ke, o5, BLUE)
card(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(0.9),
     "Closed loop: analysis → knowledge → better future response.", GREEN)

# 13 — KB derived diagram
s = new_slide(); chrome(s); title(s, "↔ Knowledge Base (KB-001 … KB-006)")
o5 = box(s, Inches(0.4), Inches(1.3), Inches(2.5), Inches(0.7), "OAS-501", fill=BLUE)
o0 = box(s, Inches(0.4), Inches(2.4), Inches(2.5), Inches(0.7), "OAS-000", fill=BLUE)
o2 = box(s, Inches(0.4), Inches(3.5), Inches(2.5), Inches(0.7), "OAS-201", fill=PURPLE)
kbs = ["KB-001 Templates","KB-002 Checklists","KB-003 Reports","KB-004 Examples","KB-005 References","KB-006 Prompts"]
ky = [1.2, 1.95, 2.7, 3.45, 4.2, 4.95]
kbox = []
for i, k in enumerate(kbs):
    kbox.append(box(s, Inches(6.6), Inches(ky[i]), Inches(4.2), Inches(0.6), k, fill=PURPLE, tcolor=WHITE, size=11))
connect(s, o5, kbox[0], PURPLE)
connect(s, o0, kbox[1], PURPLE); connect(s, o2, kbox[1], PURPLE)
connect(s, o0, kbox[2], PURPLE)
connect(s, o2, kbox[3], PURPLE)
connect(s, o0, kbox[4], PURPLE)
connect(s, o0, kbox[5], PURPLE); connect(s, o2, kbox[5], PURPLE)
box(s, Inches(0.4), Inches(5.5), Inches(5.5), Inches(0.5),
    "OAS-201 contributes: comms runbooks, QA checklist, examples & AI prompts.", fill=PANEL,
    tcolor=GREY, size=10, align=PP_ALIGN.LEFT)

# 14 — Evidence & Confidence
s = new_slide(); chrome(s); title(s, "Evidence & Confidence in MI communications")
card(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(1.0),
     "Hierarchy: Primary MI record > related > vendor > monitoring > email > user > analyst notes.", BLUE)
card(s, Inches(0.5), Inches(2.9), Inches(6.0), Inches(1.0),
     "States: classify every comms artefact — Present / Referenced / Missing / Not Applicable.", PURPLE)
card(s, Inches(6.9), Inches(1.7), Inches(5.9), Inches(1.0),
     "Confidence: rate every comms finding — High / Moderate / Low / Unknown.", GREEN)
card(s, Inches(6.9), Inches(2.9), Inches(5.9), Inches(1.0),
     "AI guardrails (OAS-000 §14): no invented evidence, no concealed uncertainty, human sign-off.", BLUE)

# 15 — Worked example table
s = new_slide(); chrome(s); title(s, "Worked example — MI000456 (Payment gateway)")
table(s, Inches(0.6), Inches(1.8), Inches(12.1), Inches(3.0),
      [["Stakeholder","Evidence","Assessment"],
       ["Operational","Bridge every 30 min","Strong"],
       ["Executive","SITREP 03:00 / 05:00 / 07:00","Adequate cadence"],
       ["Customer","Update 03:10 → next 06:45","3h35m gap during impact — weakness"],
       ["Vendor","Escalation 03:20","Effective"]])
card(s, Inches(0.6), Inches(5.1), Inches(12.1), Inches(1.0),
     "Lesson: define a max customer-comms silence window during active SEV-1.  Confidence: High.", GREEN)

# 16 — Quick start
s = new_slide(); chrome(s); title(s, "Analyst quick start — Major Incident")
bullets(s, Inches(0.7), Inches(1.8), Inches(12.0), Inches(5.0), [
    "Open OAS-201; read OAS-000 first if new.",
    "Build an Evidence Manifest; classify every comms artefact.",
    "Reconstruct the communication timeline.",
    "Assess quality & timeliness per stakeholder.",
    "Evaluate MIM handover continuity only.",
    "Assign confidence to each finding.",
    "Complete the QA checklist (KB-002).",
    "Capture Lessons for Communication; feed OAS-501 / KB.",
], size=16, gap=10)

# 17 — Takeaways
s = new_slide(); chrome(s); title(s, "Key takeaways")
card(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(1.1),
     "OAS-201 is the communications lens inside a connected ecosystem.", BLUE)
card(s, Inches(0.5), Inches(3.0), Inches(6.0), Inches(1.1),
     "All standards inherit OAS-000; none contradict.", PURPLE)
card(s, Inches(6.9), Inches(1.7), Inches(5.9), Inches(1.1),
     "The Knowledge Base is derived from the standards.", GREEN)
card(s, Inches(6.9), Inches(3.0), Inches(5.9), Inches(1.1),
     "Everything is evidence-first, confidence-rated, traceable.", BLUE)

# 18 — References
s = new_slide(); chrome(s); title(s, "References & closing")
ref = ("Standards: OAS-000 Governance · OAS-101 Incident · OAS-201 Major Incident Comms · "
       "OAS-301 Problem · OAS-401 Change · OAS-501 Knowledge.\n\n"
       "Knowledge Base: KB-001 Templates · KB-002 Checklists · KB-003 Reports · "
       "KB-004 Examples · KB-005 References · KB-006 Prompts.\n\n"
       "Release: 1.2.2   ·   Alignment: ITIL 4   ·   Reference platform: ServiceNow.")
box(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(3.2), ref, fill=PANEL, tcolor=LIGHT, size=15, align=PP_ALIGN.LEFT)
box(s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(0.9),
    "Evidence over opinion   ·   Governance before methodology   ·   Platform-agnostic, ServiceNow by reference",
    fill=PANEL, tcolor=BLUET, size=13, bold=True, align=PP_ALIGN.CENTER)

# ---------- save ----------
out = "presentation/oas-major-incident-deck.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
